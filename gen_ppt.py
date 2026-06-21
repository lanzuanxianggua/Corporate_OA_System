# -*- coding: utf-8 -*-
"""[CORPORATE OA SYSTEM] Defense PPT Generator"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = "OA_System_Defense_PPT.pptx"
IMG = "系统演示图片"

# Colors
P = RGBColor(0x66, 0x7E, 0xEA)  # primary blue
PD = RGBColor(0x4F, 0x46, 0xE5)  # dark blue
PU = RGBColor(0x76, 0x4B, 0xA2)  # purple
W = RGBColor(0xFF, 0xFF, 0xFF)
T = RGBColor(0x1E, 0x29, 0x3B)   # text
TM = RGBColor(0x47, 0x55, 0x69)  # text secondary
TS = RGBColor(0x94, 0xA3, 0xB8)  # muted
BD = RGBColor(0x0F, 0x0C, 0x29)  # dark bg
BM = RGBColor(0x1A, 0x1A, 0x2E)  # mid bg
BG = RGBColor(0xF8, 0xF9, 0xFC)  # light bg
G = RGBColor(0x05, 0x96, 0x69)   # green
O = RGBColor(0xEA, 0x58, 0x0C)   # orange
R = RGBColor(0xDC, 0x26, 0x26)   # red
TE = RGBColor(0x0D, 0x94, 0x88)  # teal

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def box(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def rbox(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def tx(s, l, t, w, h, text, sz=14, b=False, c=T, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.bold = b
    p.font.color.rgb = c; p.font.name = "Microsoft YaHei"; p.alignment = al
    return tb

def flow_num(s, l, t, n, color=P, wd=6):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(0.2), Inches(0.2))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    p = c.text_frame.paragraphs[0]; p.text = str(n)
    p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = W; p.alignment = PP_ALIGN.CENTER

def tbl(s, l, t, w, h, hd, rows, cw=None):
    rn = 1 + len(rows); cn = len(hd)
    ts = s.shapes.add_table(rn, cn, Inches(l), Inches(t), Inches(w), Inches(h))
    tb = ts.table
    if cw:
        for i, c in enumerate(cw): tb.columns[i].width = Inches(c)
    for ci, hh in enumerate(hd):
        cl = tb.cell(0, ci); cl.text = hh
        for p in cl.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = W
        cl.fill.solid(); cl.fill.fore_color.rgb = P
    for ri, rd in enumerate(rows):
        for ci, val in enumerate(rd):
            cl = tb.cell(ri+1, ci); cl.text = str(val)
            for p in cl.text_frame.paragraphs: p.font.size = Pt(9); p.font.color.rgb = T

def section_title(s, title, sub=""):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = W
    box(s, 0, 0, 13.333, 0.05, P)
    tx(s, 0.8, 0.4, 10, 0.45, title, 24, True, T)
    if sub: tx(s, 0.8, 0.8, 12, 0.3, sub, 12, False, TS)

print("Generating PPT...")

# ============ SLIDE 1: COVER ============
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = BM
box(s, 0, 7.2, 13.333, 0.05, P)
tx(s, 1.5, 2.2, 10, 1, "Enterprise OA Automation System", 38, True, W, PP_ALIGN.CENTER)
tx(s, 1.5, 3.3, 10, 0.5, "Enterprise Office Management Platform Based on Spring Boot 3.4 + Vue 3", 16, False, RGBColor(180,180,200), PP_ALIGN.CENTER)
# line placeholder - removed
box(s, 4.5, 4.1, 4.333, 0.02, P)
tx(s, 3, 4.4, 7, 1.2, "Defense Candidate: [Name]\nAdvisor: [Teacher Name]\nJune 2026", 13, False, RGBColor(150,150,175), PP_ALIGN.CENTER)
tx(s, 0.5, 6.5, 12, 0.4, "Enterprise OA Office Automation System", 11, False, TS, PP_ALIGN.CENTER)

# ============ SLIDE 2: PROJECT OVERVIEW ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Project Overview", "Enterprise Collaborative Office Platform for 200-5000 Employees")
rbox(s, 0.8, 1.4, 5.8, 3, BG)
tx(s, 1, 1.5, 5, 0.3, "Positioning", 15, True, T)
tx(s, 1, 1.85, 5.4, 2.2, "Covers 8 business domains: Administration, HR, Finance, Approval, Collaboration, Assets, Dashboard, Knowledge. 50+ functional modules that digitize daily office processes including approvals, attendance, and notifications, reducing offline paperwork and manual statistics.", 11, False, TM)

rbox(s, 6.9, 1.4, 5.8, 3, P)
tx(s, 7.1, 1.5, 5.4, 0.3, "Core Objectives", 15, True, W)
tx(s, 7.1, 1.9, 5.4, 2.2, "Configurable approval process (self-developed workflow engine)\nRBAC permission control + annotation-based auth\nData traceability (logical delete + audit fields + operation logs)\nFront-end/back-end separation with unified API\n50+ modules, easy to extend", 12, False, W)

tech = [("Backend", "Java 17 + Spring Boot 3.4 + MyBatis-Plus 3.5.9 + MySQL 8 + Redis"),
        ("Frontend", "Vue 3 + TypeScript + Vite 6 + Element Plus + Tailwind CSS + ECharts 5 + Pinia"),
        ("Tools", "Maven Multi-module + Flyway + Docker Compose + GitHub Actions CI")]
for i, (lab, cont) in enumerate(tech):
    x = 0.8 + 4.2 * i
    rbox(s, x, 4.6, 3.9, 0.85, BG)
    tx(s, x+0.2, 4.65, 3.5, 0.2, lab, 10, True, P)
    tx(s, x+0.2, 4.9, 3.5, 0.4, cont, 9, False, TM)

stats = [("50+", "Modules"), ("8", "Domains"), ("7+", "Biz Types"), ("40+", "DB Tables"), ("~2000 LoC", "Workflow Engine")]
for i, (v, l) in enumerate(stats):
    x = 0.8 + 2.5 * i
    tx(s, x, 5.65, 1.5, 0.3, v, 16, True, P, PP_ALIGN.CENTER)
    tx(s, x, 5.95, 1.5, 0.2, l, 9, False, TS, PP_ALIGN.CENTER)

# ============ SLIDE 3: HIGHLIGHT - WORKFLOW ENGINE ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Core Highlight: Self-Developed Workflow Engine", "Lightweight, Configurable, Supports 7+ Business Types with 6 Approver Strategies")

rbox(s, 0.8, 1.4, 5.8, 4.8, BG)
tx(s, 1, 1.5, 5.4, 0.3, "Why Self-Developed Instead of Flowable/Camunda?", 13, True, T)
tbl(s, 1, 1.9, 5.4, 2.0, ["Comparison", "Flowable", "Self-Developed"],
    [["Dependency", "~5 MB", "~2,000 lines of code"],
     ["DB Tables", "30+", "4 core tables"],
     ["Learning", "BPMN 2.0 required", "JSON graph, intuitive"],
     ["Approver Strategy", "Must write JavaDelegate", "Built-in 6 strategies"],
     ["Deployment", "Separate engine service", "Starts with the project"]],
    [1.2, 1.8, 2.4])

rbox(s, 1, 4.1, 5.4, 1.8, P)
tx(s, 1.1, 4.15, 5.2, 1.7,
    "Flowable: 80% features unused in OA scenarios\n"
    "  Service Task, Script Task, Receive Task, Compensate Event\n"
    "  Signal Event, Boundary Event, Business Rule Task\n\n"
    "Self-developed: Only implements OA-required capabilities\n"
    "  Approval node + conditions + 6 approver strategies + callbacks",
    10, False, W)

rbox(s, 6.9, 1.4, 5.8, 2.3, P)
tx(s, 7.1, 1.5, 5.4, 0.3, "Engine Core Components", 14, True, W)
engine_parts = [("1. Process Definition", "JSON graph (nodes + edges + conditions)"),
    ("2. Process Instance", "Linked business document + snapshot isolation"),
    ("3. Approval Task", "Todo/Done/Countersign/Orsign/Delegation"),
    ("4. Runtime Engine", "Graph parsing + condition evaluation + routing"),
    ("5. Business Callback", "Spring Event async notification")]
for i, (t, d) in enumerate(engine_parts):
    flow_num(s, 7.1, 1.9 + i*0.42, i+1, W)
    tx(s, 7.35, 1.9 + i*0.42, 5, 0.35, f" {t} -- {d}", 9, False, W)

rbox(s, 6.9, 3.9, 5.8, 2.3, BG)
tx(s, 7.1, 4, 5.4, 0.3, "Supported Business Types", 13, True, T)
biz = [("Leave Approve", P), ("Biz Trip", TE), ("Expense", O), ("Purchase", PU), ("Overtime", G), ("Outing", TM), ("Loan", R)]
for i, (n, c) in enumerate(biz):
    x = 7.1 + (i%4)*1.35; y = 4.35 + (i//4)*0.4
    rbox(s, x, y, 1.2, 0.32, c)
    tx(s, x, y+0.02, 1.2, 0.28, n, 8, True, W, PP_ALIGN.CENTER)

tx(s, 0.8, 6.4, 5.8, 0.4, "Snapshot: Process starts with a copy of definition, later changes don't affect running instances.", 9, False, TM)
tx(s, 6.9, 6.4, 5.8, 0.4, "Lock: Redis SET NX EX 10 prevents concurrent conflicts in countersign/orsign.", 9, False, TM)

# ============ SLIDE 4: JSON GRAPH & CONDITION ROUTING ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "JSON Graph Structure & Condition Routing", "Schema v3, 4-Dimension Routing, 6 Approver Strategies")

rbox(s, 0.8, 1.4, 6.8, 4.5, BD)
tx(s, 1, 1.5, 6.4, 0.3, "Amount-Tiered Approval JSON (Simplified)", 12, True, W)
code = """{
  "schemaVersion": 3,
  "nodes": [
    { "nodeId": "start", "nodeType": "start" },
    { "nodeId": "manager", "nodeType": "approval",
      "nodeName": "Department Manager",
      "assigneeType": "dept_manager" },
    { "nodeId": "finance", "nodeType": "approval",
      "nodeName": "Finance Review",
      "assigneeType": "role_global",
      "assigneeValue": "FINANCE",
      "conditions": [{"field":"amount",
        "operator":">","value":5000}] },
    { "nodeId": "director", "nodeType": "approval",
      "nodeName": "Director",
      "assigneeType": "role_global",
      "assigneeValue": "DIRECTOR",
      "conditions": [{"field":"amount",
        "operator":">","value":50000}] },
    { "nodeId": "end", "nodeType": "end" }
  ],
  "edges": [...]
}"""
tx(s, 1, 1.85, 6.4, 3.8, code, 9, False, RGBColor(200,210,220))

rbox(s, 7.9, 1.4, 4.8, 1.8, BG)
tx(s, 8.1, 1.5, 4.4, 0.3, "4-Dimension Routing", 13, True, T)
tbl(s, 8.1, 1.85, 4.4, 1.2, ["Dimension", "Field", "Example"],
    [["Amount", "amount", "amount > 5000"],
     ["Days", "days", "days >= 7"],
     ["Hours", "hours", "hours > 8"],
     ["Level", "role_level", "initiator_level_match"]],
    [1.2, 1.0, 2.2])

rbox(s, 7.9, 3.4, 4.8, 2.5, BG)
tx(s, 8.1, 3.5, 4.4, 0.3, "6 Approver Strategies", 13, True, T)
tbl(s, 8.1, 3.85, 4.4, 1.8, ["Strategy", "Scope", "Use Case"],
    [["dept_manager", "Direct superior", "Leave, Outing"],
     ["role", "Same dept role", "Dept internal"],
     ["role_global", "Company-wide role", "Finance/HR"],
     ["specific", "Specific employee", "Fixed approver"],
     ["role_chain", "Role chain top-down", "Tiered approval"],
     ["initiator_level_match", "Skip based on level", "Executive skip"]],
    [1.4, 1.5, 1.5])

rbox(s, 0.8, 6.1, 12, 0.5, BG)
tx(s, 1, 6.15, 11.5, 0.35, "Operators: >, >=, <, <=, ==, !=, contains, in, starts_with | Supports && and || compound expressions", 10, False, TM)

# ============ SLIDE 5: CUSTOM DEFINITION (WITH SCREENSHOT) ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Custom Process Definition", "Visual Designer -> JSON Graph -> Engine Execution, Zero-Code Configuration")

if os.path.exists(f"{IMG}/流程自定义.png"):
    s.shapes.add_picture(f"{IMG}/流程自定义.png", Inches(0.3), Inches(1.3), Inches(9.0), Inches(4.8))

rbox(s, 9.5, 1.3, 3.5, 1.4, P)
tx(s, 9.7, 1.35, 3.1, 0.25, "One-Click Templates", 12, True, W)
tx(s, 9.7, 1.65, 3.1, 0.8, "Standard: Dept Manager -> HR Review\nAmount-Tiered: Auto-route by amount\nDuration-Tiered: Auto-route by days", 9, False, W)

rbox(s, 9.5, 2.8, 3.5, 1.1, BG)
tx(s, 9.7, 2.85, 3.1, 0.25, "Online Validation", 12, True, T)
tx(s, 9.7, 3.1, 3.1, 0.6, "Auto-detects 11 rules before saving:\nCycle detection, connectivity check,\nnode ID uniqueness, approver config", 9, False, TM)

rbox(s, 9.5, 4.0, 3.5, 1.0, BG)
tx(s, 9.7, 4.05, 3.1, 0.25, "Version Management", 12, True, T)
tx(s, 9.7, 4.3, 3.1, 0.5, "Auto version+1 per edit\nOld version disabled, new activated\nRunning instances unaffected (snapshot)", 9, False, TM)

rbox(s, 9.5, 5.1, 3.5, 1.1, RGBColor(255,247,237))
tx(s, 9.7, 5.15, 3.1, 0.25, "Designer Features", 11, True, O)
tx(s, 9.7, 5.4, 3.1, 0.6, "Step reorder/delete, approver search,\ncountersign/orsign toggle,\ntimeout config (hours)", 9, False, TM)

rbox(s, 0.3, 6.3, 12.7, 0.55, BG)
tx(s, 0.5, 6.35, 12.3, 0.4, "Flow: Admin opens designer -> Add approval steps / set tier rules -> Save -> Backend validates JSON -> Version+1 activated -> Employee submits -> Engine routes -> Create todo -> WebSocket notify", 10, False, T)

# ============ SLIDE 6: COUNTERSIGN / ORSIGN / TIMEOUT ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Countersign, Orsign & Timeout Escalation", "Multi-person Approval + Auto-escalation + Redis Distributed Lock")

for i, (tt, body, c) in enumerate([
    ("Countersign", "ALL approvers must approve.\n\nA approves -> Check B,C -> Wait\nB approves -> Check C -> Wait\nC approves -> ALL done -> Next node", P),
    ("Orsign", "ANY approver can approve.\n\nA approves -> Cancel B,C tasks\n-> Advance to next node\n-> B,C see \"handled by others\"", TE),
    ("Timeout Escalation", "If approver is late, escalate.\n\ntimeoutHours: 48\ntimeoutAction: \"escalate\"\nescalateTo: GM", O)]):
    x = 0.8 + 4.2 * i
    rbox(s, x, 1.5, 3.9, 3.2, c)
    tx(s, x+0.2, 1.6, 3.5, 0.4, tt, 14, True, W)
    tx(s, x+0.2, 2.05, 3.5, 2.4, body, 11, False, W)

rbox(s, 0.8, 5.0, 12, 0.8, BG)
tx(s, 1, 5.05, 11.5, 0.6, "Redis Distributed Lock: SET key uuid NX EX 10 prevents concurrent escalation. Parent-child task structure (parentTask + childTask) manages multi-person scenarios. Countersign checks sibling count, orsign cancels remaining siblings.", 10, False, T)

rbox(s, 0.8, 6.0, 5.8, 0.6, BG)
tx(s, 1, 6.05, 5.4, 0.4, "Spring Event callback: On approval/rejection, publish event -> listener auto-handles balance deduction and attendance marking.", 9, False, TM)
rbox(s, 6.9, 6.0, 5.8, 0.6, BG)
tx(s, 7.1, 6.05, 5.4, 0.4, "Timeout escalation auto-escalates to higher-level approver when config timeout is reached. Supports auto-approve and auto-reject.", 9, False, TM)

# ============ SLIDE 7: REDIS 6 SCENARIOS ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Redis: 6 Usage Scenarios", "One Cache Middleware, Multiple Core Roles")

items = [
    ("JWT Token Mgmt", "Token stored in Redis (EX 7200s)\nAdmin force-logout = delete key\nJWT stateless + Redis controllable", P),
    ("IP Rate Limiting", "INCR rate:login:ip -> EXPIRE 60s\nMax 5 attempts/min per IP\nAtomic operation, no manual cleanup", O),
    ("Online Users", "online:user:{empId} 30min heartbeat\nAdmin can view/force-offline users\nKEYS scan to get all online users", TE),
    ("Captcha Storage", "captcha:{uuid} -> code EX 300s\nDeleted immediately after verification\nOne-time use prevents replay", PU),
    ("Token Refresh", "Frontend detects expiry 5min early\nAuto-call /refresh-token\nNew token extends Redis 2h", P),
    ("Distributed Lock", "SET key uuid NX EX 10\nGET -> compare uuid -> DEL\nPrevents concurrent task handling", O),
]
for i, (tt, desc, c) in enumerate(items):
    col, row = i % 3, i // 3
    x = 0.8 + 4.2 * col; y = 1.5 + 2.8 * row
    rbox(s, x, y, 3.9, 2.5, BG)
    box(s, x, y, 3.9, 0.05, c)
    tx(s, x+0.15, y+0.15, 3.6, 0.3, tt, 13, True, c)
    tx(s, x+0.15, y+0.5, 3.6, 1.8, desc, 10, False, TM)

rbox(s, 0.8, 7.0, 12, 0.3, P)
tx(s, 1, 7.0, 11.5, 0.25, "Strategy: Leverage Redis INCR atomicity + EXPIRE auto-expiry for lightweight rate limiting & caching without additional middleware.", 10, False, W)

# ============ SLIDE 8: JWT AUTH & PERMISSION ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "JWT Authentication & Permission System", "JWT + Redis Dual Verification + RBAC + Annotation-based Auth")

rbox(s, 0.8, 1.4, 5.8, 3.0, BG)
tx(s, 1, 1.5, 5.4, 0.3, "Auth Flow", 14, True, T)
steps = [("Login", "Generate JWT (HMAC-SHA256)"), ("Redis", "token: + roles: keys EX 7200s"),
         ("Frontend", "Store in localStorage, Axios interceptor"), ("Interceptor", "AuthInterceptor parses JWT + checks Redis"),
         ("Force Logout", "Admin deletes Redis key, instant effect")]
for i, (tt, d) in enumerate(steps):
    y = 1.95 + i * 0.42
    rbox(s, 1.1, y, 1.2, 0.35, P)
    tx(s, 1.15, y+0.03, 1.1, 0.3, tt, 9, True, W, PP_ALIGN.CENTER)
    tx(s, 2.45, y+0.05, 3.8, 0.3, d, 10, False, TM)
    if i < len(steps)-1: tx(s, 1.55, y+0.32, 0.5, 0.2, "v", 7, False, TS, PP_ALIGN.CENTER)

rbox(s, 6.9, 1.4, 5.8, 1.5, BG)
tx(s, 7.1, 1.5, 5.4, 0.3, "RBAC Model", 14, True, T)
entities = [("User", P), ("Role", TE), ("Menu", O), ("Permission", PU)]
for i, (n, c) in enumerate(entities):
    x = 7.1 + i * 1.3
    rbox(s, x, 1.95, 1.0, 0.35, c)
    tx(s, x, 2.0, 1.0, 0.3, n, 9, True, W, PP_ALIGN.CENTER)
    if i < len(entities)-1: tx(s, x+1.0, 2.0, 0.3, 0.3, "->", 11, True, TS, PP_ALIGN.CENTER)

rbox(s, 6.9, 3.1, 5.8, 1.0, BD)
tx(s, 7.1, 3.15, 5.4, 0.25, "Annotation-based Auth", 11, True, W)
tx(s, 7.1, 3.4, 5.4, 0.5, '@RequirePermission("hr:leave:create")\n@RequireAdmin\n@RequireRole({"ADMIN", "DEPT_MANAGER"})', 9, False, RGBColor(200,210,220))

rbox(s, 0.8, 4.6, 12, 0.8, RGBColor(255,247,237))
tx(s, 1, 4.65, 11.5, 0.6, "Q: Why store JWT in Redis if it's already stateless?\nA: JWT can't be revoked once issued. Redis enables server-side invalidation -- admin \"force-logout\" is just a DEL command. Password changes also auto-delete old tokens.", 10, False, T)

# ============ SLIDE 9: LOGIN ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Login Implementation", "Animated Characters + Captcha + BCrypt + JWT")

if os.path.exists(f"{IMG}/登录页.png"):
    s.shapes.add_picture(f"{IMG}/登录页.png", Inches(0.3), Inches(1.3), Inches(7.8), Inches(4.5))

rbox(s, 8.3, 1.3, 4.7, 4.5, BG)
tx(s, 8.5, 1.4, 4.3, 0.3, "Tech Highlights", 14, True, T)
items = [
    "Eye tracking: Pupil/EyeBall Vue components, mousemove + computed offset",
    "Blink: setInterval random 3-7s, blink duration 150ms",
    "Password toggle: Eye icon switches type=text/password",
    "Captcha: Base64 image, Redis 5min expiry, one-time use",
    "BCrypt: One-way hash, verification via checkpw()",
    "Repeat prevention: loading lock + disabled button",
]
for i, it in enumerate(items):
    flow_num(s, 8.5, 1.85 + i*0.45, i+1, P)
    tx(s, 8.75, 1.85 + i*0.45, 4.2, 0.35, it, 9, False, TM)

# ============ SLIDE 10: SCREENSHOT GRID ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "System Screenshots", "Feature Showcase")

imgs_top = [("工作台.png", 0.3, 1.3, 6.5, 3.2), ("员工管理.png", 6.9, 1.3, 6.0, 3.2)]
imgs_bot = [("考勤打卡.png", 0.3, 4.6, 4.3, 2.5), ("消息中心.png", 4.7, 4.6, 4.3, 2.5), ("请假申请.png", 9.1, 4.6, 4.0, 2.5)]
for fname, x, y, w, h in imgs_top + imgs_bot:
    path = f"{IMG}/{fname}"
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

# ============ SLIDE 11: DATABASE & FLYWAY ============
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Database Design & Flyway Migration", "Domain-based Design + Audit Fields + Versioned Migration")

rbox(s, 0.8, 1.4, 5.8, 2.3, BG)
tx(s, 1, 1.5, 5.4, 0.3, "Domain Classification", 14, True, T)
tbl(s, 1, 1.9, 5.4, 1.5, ["Domain", "Prefix", "Tables"],
    [["System", "sys_", "Employee, Dept, Role, Menu"],
     ["OA Business", "oa_", "Attendance, Leave, Message"],
     ["Workflow", "wf_", "Definition, Instance, Task"],
     ["Extension", "hr_/fin_", "Performance, Payroll"]],
    [1.5, 1.0, 2.9])

rbox(s, 0.8, 3.9, 5.8, 1.0, BG)
tx(s, 1, 3.95, 5.4, 0.25, "Audit Fields", 11, True, T)
tx(s, 1, 4.2, 5.4, 0.5, "Every table: create_time, update_time, create_by, update_by, del_flag (logical delete), version (optimistic lock). Auto-filled by MyMetaObjectHandler.", 9, False, TM)

rbox(s, 6.9, 1.4, 5.8, 2.3, BG)
tx(s, 7.1, 1.5, 5.4, 0.3, "Flyway Migration", 14, True, T)
tx(s, 7.1, 1.9, 5.4, 0.8, "V1011__workflow_graph_seeds.sql\nV1012__tiered_workflow_approval_chains.sql\nV1013__workflow_schema_v3_custom_runtime.sql", 9, False, TM)
tx(s, 7.1, 2.7, 5.4, 0.5, "Auto-executes on startup\nTracked via flyway_schema_history\nReset: DROP TABLE flyway_schema_history", 9, False, TM)

rbox(s, 6.9, 3.9, 5.8, 1.0, BG)
tx(s, 7.1, 3.95, 5.4, 0.25, "Index Strategy", 11, True, T)
tx(s, 7.1, 4.2, 5.4, 0.5, "Foreign keys (dept_id, assignee) and status fields indexed\nComposite (emp_id, work_date) for attendance queries\nUnique emp_code and phone for data integrity", 9, False, TM)

# ============ SLIDE 12: HIGHLIGHTS ============
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = BM
tx(s, 0.8, 0.3, 10, 0.5, "Project Highlights", 24, True, W)
tx(s, 0.8, 0.75, 12, 0.3, "Technical Depth + Business Coverage + Engineering Excellence", 12, False, RGBColor(180,180,200))

hl = [
    ("Self-Developed Workflow", "JSON graph, 4-dim routing, 6 approver strategies. 7+ business types with 2000 LoC vs Flowable's 5MB.", P),
    ("Redis in 6 Scenarios", "Token mgmt, rate limiting, online users, captcha, token refresh, distributed lock. One cache, 6 roles.", TE),
    ("JWT + Redis Dual Auth", "Stateless JWT + server-side Redis revocation. Annotation-based @RequirePermission + RBAC.", O),
    ("Spring Event Decoupling", "Approval complete -> Event -> Listener auto-handles balance deduct & attendance marking. Zero coupling.", PU),
    ("ECharts Dashboard", "11 charts (trend, heatmap, radar, scatter, gauge). Dual Y-axis, theme switching, responsive.", G),
    ("Unified Excel Export", "All exports follow: FileName_YYYYMMDD_HHmmss.xlsx. Auto column width, status code to Chinese.", R),
]
for i, (tt, d, c) in enumerate(hl):
    col, row = i % 3, i // 3
    x = 0.8 + 4.2 * col; y = 1.2 + 2.8 * row
    rbox(s, x, y, 3.9, 2.6, BG)
    box(s, x, y, 3.9, 0.05, c)
    tx(s, x+0.15, y+0.15, 3.6, 0.35, tt, 13, True, c)
    tx(s, x+0.15, y+0.55, 3.6, 1.8, d, 10, False, TM)

tx(s, 0.5, 6.8, 12.3, 0.3, "Comprehensive Coverage  |  Complete Permission System  |  Data Traceability  |  Mobile Accessible  |  Front-Back Separation", 10, False, RGBColor(200,200,220), PP_ALIGN.CENTER)

# ============ SLIDE 13: THANK YOU ============
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = BM
box(s, 0, 7.2, 13.333, 0.05, P)
tx(s, 1, 2.5, 11, 1, "Thank You", 48, True, W, PP_ALIGN.CENTER)
tx(s, 1, 3.6, 11, 0.5, "Questions & Discussion Welcome", 20, False, RGBColor(180,180,200), PP_ALIGN.CENTER)
tx(s, 2, 4.8, 9.3, 1, "Enterprise OA Office Automation System\nBased on Spring Boot 3.4 + Vue 3\nDefense Candidate: [Name] | Advisor: [Teacher] | June 2026", 12, False, RGBColor(150,150,175), PP_ALIGN.CENTER)

print(f"Saving PPT ({len(prs.slides)} slides)...")
prs.save(OUTPUT)
print(f"Success: {OUTPUT}")
